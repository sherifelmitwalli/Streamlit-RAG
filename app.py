import streamlit as st
import os
import logging
import time
from io import BytesIO
from typing import List, Optional, Dict, Any, Tuple
import zipfile  # For handling ZIP files
import fitz  # PyMuPDF for PDF extraction
import pandas as pd
import numpy as np
from pptx import Presentation
from sklearn.metrics.pairwise import cosine_similarity
import openai
import json
import re

# -----------------------------
# Logging & Constants
# -----------------------------
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(name)s - %(message)s",
    handlers=[
        logging.StreamHandler(),
        logging.FileHandler("app.log", mode="a")
    ]
)
logger = logging.getLogger(__name__)

MAX_FILE_SIZE = 50 * 1024 * 1024  # 50 MB
CHUNK_SIZE = 500                  # Maximum characters per text chunk
OVERLAP_SIZE = int(CHUNK_SIZE * 0.2)  # 20% overlap for better context
MAX_RETRIES = 3                   # Maximum number of API call retries
RETRY_DELAY = 5                   # Delay between retries in seconds
EMBEDDING_MODEL = "text-embedding-ada-002"
TOP_K = 5                         # Increased from 3 to 5 for better coverage

# -----------------------------
# OpenAI & Secrets Setup
# -----------------------------
try:
    openai.api_key = st.secrets["OPENAI_API_KEY"]
    MODEL_NAME = st.secrets["MODEL"]
    if not openai.api_key:
        raise ValueError("OPENAI_API_KEY not configured in Streamlit secrets")
    if not MODEL_NAME:
        raise ValueError("MODEL not configured in Streamlit secrets")
except (KeyError, ValueError) as e:
    error_msg = str(e)
    logger.error(f"Configuration error: {error_msg}")
    st.error("Configuration Error: Missing required secrets.")
    st.stop()

# Validate API key by making a test request
try:
    response = openai.Embedding.create(
        input="test",
        model=EMBEDDING_MODEL
    )
    if response and response.get('data'):
        logger.info("Successfully connected to OpenAI API")
    else:
        raise ValueError("Invalid API response")
except Exception as e:
    error_msg = f"Failed to connect to OpenAI API: {str(e)}"
    logger.error(error_msg)
    st.error(error_msg)
    st.stop()

# -----------------------------
# Custom Styling & Sidebar
# -----------------------------
st.markdown('''
<style>
div.stButton > button:first-child {
    background-color: #4CAF50;
    color: white;
    font-size: 16px;
    padding: 10px 24px;
    border-radius: 8px;
}
</style>
''', unsafe_allow_html=True)

st.sidebar.header("About")
st.sidebar.info("This app uses AI-assisted Retrieval-Augmented Generation (RAG) to answer questions based on uploaded files or folders.")
st.sidebar.markdown("[View Documentation](https://example.com)")

st.title("Agentic RAG Chatbot")
st.subheader("Upload one or more files or a folder (as a ZIP file) and ask questions based on their content.")

# -----------------------------
# Helper: Clear Cache
# -----------------------------
def clear_cache() -> None:
    try:
        st.cache_data.clear()
        logger.info("Cache cleared successfully")
    except Exception as e:
        logger.error(f"Error clearing cache: {e}")

# -----------------------------
# Helper: Clean Text
# -----------------------------
def clean_text(text: str) -> str:
    """
    Enhanced text cleaning:
    - Normalizes whitespace
    - Removes hyphenation artifacts
    - Preserves sentence boundaries
    - Handles common PDF artifacts
    """
    # Remove PDF artifacts and normalize whitespace
    text = re.sub(r'\f', ' ', text)  # Form feed
    text = re.sub(r'(?<=[a-z])-\s*\n\s*(?=[a-z])', '', text)  # Remove hyphenation
    text = re.sub(r'\s*\n\s*', ' ', text)  # Normalize line breaks
    text = re.sub(r'\s+', ' ', text)  # Collapse multiple spaces
    text = re.sub(r'(?<=[.!?])\s+', '\n', text)  # Preserve sentence boundaries
    return text.strip()

# -----------------------------
# Text Processing Functions
# -----------------------------
def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = OVERLAP_SIZE) -> List[str]:
    """
    Improved chunking strategy that preserves sentences and ensures proper context.
    """
    sentences = re.split(r'(?<=[.!?])\s+', text)
    chunks = []
    current_chunk = []
    current_length = 0
    last_sentences = []  # Keep track of last sentences for overlap
    
    for sentence in sentences:
        sentence_length = len(sentence)
        
        # If this sentence alone exceeds chunk size, split it
        if sentence_length > chunk_size:
            if current_chunk:  # First store any accumulated chunk
                chunks.append(' '.join(current_chunk))
                last_sentences = current_chunk[-2:] if len(current_chunk) > 2 else current_chunk
                current_chunk = []
                current_length = 0
            
            # Split long sentence while preserving words
            words = sentence.split()
            current_words = []
            current_word_count = 0
            for word in words:
                if current_word_count + len(word) > chunk_size:
                    if current_words:
                        chunks.append(' '.join(current_words))
                    current_words = [word]
                    current_word_count = len(word)
                else:
                    current_words.append(word)
                    current_word_count += len(word) + 1
            if current_words:
                chunks.append(' '.join(current_words))
            continue
            
        # Normal sentence processing
        if current_length + sentence_length <= chunk_size:
            current_chunk.append(sentence)
            current_length += sentence_length + 1  # +1 for space
        else:
            if current_chunk:
                chunks.append(' '.join(current_chunk))
                # Keep last sentences for overlap
                last_sentences = current_chunk[-2:] if len(current_chunk) > 2 else current_chunk
                # Start new chunk with overlap
                current_chunk = last_sentences + [sentence]
                current_length = sum(len(s) + 1 for s in current_chunk)
            else:
                # If we have no current chunk but the sentence is small enough
                current_chunk = [sentence]
                current_length = sentence_length
                
    # Don't forget the last chunk
    if current_chunk:
        chunks.append(' '.join(current_chunk))
    
    return chunks

def process_pdf_blocks(doc: fitz.Document, file_name: str) -> List[Dict[str, Any]]:
    """
    Enhanced PDF processing using PyMuPDF's block extraction for better layout handling.
    """
    chunks = []
    for page_num in range(len(doc)):
        page = doc[page_num]
        # Extract text blocks with layout information
        blocks = page.get_text("dict")["blocks"]
        page_text = []
        
        for block in blocks:
            if block.get("type") == 0:  # Type 0 is text
                block_text = []
                for line in block.get("lines", []):
                    line_text = []
                    for span in line.get("spans", []):
                        text = span.get("text", "").strip()
                        if text:
                            line_text.append(text)
                    if line_text:
                        block_text.append(" ".join(line_text))
                if block_text:
                    page_text.append(" ".join(block_text))
        
        if page_text:
            text = clean_text("\n".join(page_text))
            # Split into smaller chunks while preserving context
            text_chunks = chunk_text(text)
            for chunk in text_chunks:
                chunks.append({
                    "text": f"File: {file_name} | Page: {page_num + 1}\n{chunk}",
                    "source": {"file": file_name, "page": page_num + 1}
                })
    
    return chunks

def process_file_bytes(file_bytes: BytesIO, file_name: str) -> Optional[List[Dict[str, Any]]]:
    """
    Enhanced file processing with improved text extraction and chunking.
    """
    ext = file_name.split('.')[-1].lower()
    try:
        if ext == "txt":
            file_bytes.seek(0)
            text = file_bytes.read().decode("utf-8")
            text = clean_text(text)
            chunks = []
            for chunk in chunk_text(text):
                chunks.append({
                    "text": f"File: {file_name}\n{chunk}",
                    "source": {"file": file_name}
                })
            return chunks
        elif ext == "pdf":
            file_bytes.seek(0)
            pdf_data = file_bytes.read()
            doc = fitz.open(stream=pdf_data, filetype="pdf")
            chunks = process_pdf_blocks(doc, file_name)
            doc.close()
            return chunks
        elif ext == "csv":
            file_bytes.seek(0)
            df = pd.read_csv(file_bytes)
            text = df.to_string()
            text = clean_text(text)
            chunks = []
            for chunk in chunk_text(text):
                chunks.append({
                    "text": f"File: {file_name}\n{chunk}",
                    "source": {"file": file_name}
                })
            return chunks
        elif ext == "pptx":
            file_bytes.seek(0)
            presentation = Presentation(file_bytes)
            text = "\n".join(
                shape.text for slide in presentation.slides 
                for shape in slide.shapes if hasattr(shape, "text")
            )
            text = clean_text(text)
            chunks = []
            for chunk in chunk_text(text):
                chunks.append({
                    "text": f"File: {file_name}\n{chunk}",
                    "source": {"file": file_name}
                })
            return chunks
        elif ext == "json":
            return process_json_file(file_bytes, file_name)
        else:
            st.error(f"Unsupported file extension: {ext} for file {file_name}")
            return None
    except Exception as e:
        logger.error(f"Error processing file {file_name}: {str(e)}", exc_info=True)
        st.error(f"Error processing file {file_name}. Error: {str(e)}")
        return None

def process_json_file(file_bytes: BytesIO, file_name: str) -> List[Dict[str, Any]]:
    """
    Process a JSON file and convert it to text chunks with metadata.
    """
    chunks = []
    try:
        file_bytes.seek(0)
        data = json.load(file_bytes)
    except json.JSONDecodeError as jde:
        error_msg = f"JSON decoding error in file {file_name}: {str(jde)}"
        logger.error(error_msg, exc_info=True)
        st.error(error_msg)
        return []
    except Exception as e:
        error_msg = f"Unexpected error reading JSON file {file_name}: {str(e)}"
        logger.error(error_msg, exc_info=True)
        st.error(error_msg)
        return []
    
    if isinstance(data, dict):
        for key, value in data.items():
            if isinstance(key, str) and re.fullmatch(r'page_\d+', key.lower()):
                page_num = key.split("_")[-1]
                text = json.dumps(value, indent=2) if isinstance(value, (dict, list)) else str(value)
                chunks.extend([{
                    "text": f"File: {file_name} | Page: {page_num}\n{chunk}",
                    "source": {"file": file_name, "page": page_num}
                } for chunk in chunk_text(text)])
            else:
                text = json.dumps(value, indent=2) if isinstance(value, (dict, list)) else str(value)
                chunks.extend([{
                    "text": f"File: {file_name} | Key: {key}\n{chunk}",
                    "source": {"file": file_name, "key": key}
                } for chunk in chunk_text(text)])
    elif isinstance(data, list):
        for index, item in enumerate(data, start=1):
            text = json.dumps(item, indent=2)
            chunks.extend([{
                "text": f"File: {file_name} | Item: {index}\n{chunk}",
                "source": {"file": file_name, "item": index}
            } for chunk in chunk_text(text)])
    else:
        text = json.dumps(data, indent=2)
        chunks.extend([{
            "text": f"File: {file_name}\n{chunk}",
            "source": {"file": file_name}
        } for chunk in chunk_text(text)])
    
    return chunks

def process_file(uploaded_file: st.runtime.uploaded_file_manager.UploadedFile) -> Optional[List[Dict[str, Any]]]:
    """
    Process uploaded file and handle any errors.
    """
    try:
        file_bytes = BytesIO(uploaded_file.read())
        return process_file_bytes(file_bytes, uploaded_file.name)
    except Exception as e:
        logger.error(f"Error processing file {uploaded_file.name}: {str(e)}", exc_info=True)
        st.error(f"Error processing file {uploaded_file.name}. Error: {str(e)}")
        return None
    finally:
        file_bytes.close()

# -----------------------------
# Embedding & Search Functions
# -----------------------------
@st.cache_data(show_spinner=False, ttl=3600)
def generate_embeddings(text_chunks: List[str]) -> np.ndarray:
    """
    Generate embeddings for text chunks with improved error handling and rate limiting.
    """
    try:
        batch_size = 20
        all_embeddings = []
        for i in range(0, len(text_chunks), batch_size):
            batch = [chunk.strip() for chunk in text_chunks[i:i + batch_size] if chunk.strip()]
            if not batch:
                continue
            
            for attempt in range(MAX_RETRIES):
                try:
                    response = openai.Embedding.create(
                        input=batch,
                        model=EMBEDDING_MODEL
                    )
                    batch_embeddings = [data['embedding'] for data in response['data']]
                    all_embeddings.extend(batch_embeddings)
                    if i + batch_size < len(text_chunks):
                        time.sleep(0.5)  # Rate limiting
                    break
                except openai.error.RateLimitError:
                    if attempt < MAX_RETRIES - 1:
                        st.warning(f"Rate limit reached. Waiting {RETRY_DELAY} seconds...")
                        time.sleep(RETRY_DELAY * (attempt + 1))  # Exponential backoff
                    else:
                        raise
        
        if not all_embeddings:
            raise ValueError("No valid embeddings were generated")
        return np.array(all_embeddings)
    except Exception as e:
        logger.error(f"Error generating embeddings: {str(e)}", exc_info=True)
        st.error(f"Failed to generate embeddings: {str(e)}")
        return np.array([])

def get_chat_response(messages: List[Dict[str, str]], retries: int = MAX_RETRIES, delay: int = RETRY_DELAY) -> Optional[str]:
    """
    Get response from OpenAI Chat API with retries and error handling.
    """
    for attempt in range(retries):
        try:
            response = openai.ChatCompletion.create(
                model=MODEL_NAME,
                messages=messages,
                temperature=0.7,
                max_tokens=1000
            )
            return response.choices[0].message["content"]
        except openai.error.RateLimitError:
            if attempt < retries - 1:
                st.warning(f"Rate limit exceeded. Retrying in {delay} seconds...")
                time.sleep(delay * (attempt + 1))  # Exponential backoff
            else:
                raise
        except Exception as e:
            logger.error(f"OpenAI API error: {e}", exc_info=True)
            st.error("An error occurred while generating the response. Please try again later.")
            return None
    return None

def find_relevant_context_indices(query: str, text_chunks: List[str], embeddings: np.ndarray, top_k: int = TOP_K) -> List[int]:
    """
    Find relevant context using improved similarity search.
    """
    try:
        # Get query embedding
        response = openai.Embedding.create(
            input=[query],
            model=EMBEDDING_MODEL
        )
        query_embedding = np.array(response['data'][0]['embedding']).reshape(1, -1)
        
        # Calculate similarities
        similarities = cosine_similarity(query_embedding, embeddings)[0]
        
        # Get top k indices with minimum similarity threshold
        min_similarity = 0.3  # Minimum similarity threshold
        mask = similarities >= min_similarity
        filtered_indices = np.where(mask)[0]
        
        if len(filtered_indices) > 0:
            # Sort by similarity and take top_k
            top_indices = filtered_indices[np.argsort(similarities[filtered_indices])[-top_k:][::-1]]
            return list(top_indices)
        else:
            # Fallback to standard top_k if no chunks meet threshold
            top_indices = similarities.argsort()[-top_k:][::-1]
            return list(top_indices)
            
    except Exception as e:
        logger.error(f"Error finding relevant context: {e}", exc_info=True)
        st.error("Failed to retrieve relevant context. Please try again later.")
        return []

def extract_snippets(text: str, search_terms: List[str], context: int = 100) -> List[str]:
    """
    Extract multiple snippets for a list of search terms with improved context handling.
    """
    # Remove metadata line for matching
    # Get the metadata line first
    metadata_match = re.match(r'^File:.*?\n', text)
    metadata_length = len(metadata_match.group(0)) if metadata_match else 0
    content_text = text[metadata_length:]
    
    snippets = []
    text_lower = content_text.lower()
    
    for term in search_terms:
        term_lower = term.lower()
        matches = list(re.finditer(rf'\b{re.escape(term_lower)}\b', text_lower))
        
        for match in matches:
            # Adjust start and end positions to account for metadata
            start = max(metadata_length, metadata_length + match.start() - context)
            end = min(len(text), metadata_length + match.end() + context)
            
            # Expand to word boundaries
            if start > 0:
                start = text.rfind(' ', 0, start) + 1
            if end < len(text):
                end = text.find(' ', end)
                if end == -1:
                    end = len(text)
            
            snippet = text[start:end].strip()
            if snippet:
                # Add ellipsis if snippet is truncated
                if start > 0:
                    snippet = "..." + snippet
                if end < len(text):
                    snippet = snippet + "..."
                snippets.append(snippet)
    
    return snippets

def flexible_match(text: str, term: str) -> Tuple[bool, float]:
    """
    Enhanced matching with improved word boundaries and sensitivity.
    """
    # Clean and normalize the input text
    text_clean = text.lower()
    # Remove File/Page prefix lines for matching
    text_clean = re.sub(r'^File:.*?\|.*?\n', '', text_clean)
    text_clean = clean_text(text_clean)
    
    term_clean = term.lower()
    
    # First try exact phrase matching
    if f" {term_clean} " in f" {text_clean} ":
        return True, 1.0
    
    # Check for term at start or end of text
    if text_clean.startswith(term_clean + " ") or text_clean.endswith(" " + term_clean):
        return True, 1.0
        
    # Then try word-by-word matching
    words = term_clean.split()
    if len(words) == 1:
        # For single words, use word boundary matching
        pattern = r'\b' + re.escape(term_clean) + r'\b'
        if re.search(pattern, text_clean, re.IGNORECASE):
            return True, 1.0
    else:
        # For phrases, check each word
        word_matches = 0
        for word in words:
            pattern = r'\b' + re.escape(word) + r'\b'
            if re.search(pattern, text_clean, re.IGNORECASE):
                word_matches += 1
        match_score = word_matches / len(words)
        if match_score == 1.0:
            return True, 1.0
        elif match_score > 0.5:  # More than half the words match
            return False, match_score
    
    # Last resort: check for fuzzy matches on longer terms
    if len(term_clean) > 4:
        for text_word in text_clean.split():
            if len(text_word) > 4:
                if text_word.startswith(term_clean[:4]) or text_word.endswith(term_clean[-4:]):
                    return False, 0.8
    
    return False, 0.0

def extract_exact_mentions(chunks: List[Dict[str, Any]], search_term: str, min_score: float = 0.8) -> List[Dict[str, Any]]:
    """
    Enhanced exact mention extraction with better context and scoring.
    """
    results = []
    search_words = search_term.lower().split()
    
    for chunk in chunks:
        original_text = chunk.get("text", "")
        exact_match, match_score = flexible_match(original_text, search_term)
        
        if match_score >= min_score:
            snippets = extract_snippets(original_text, search_words)
            if snippets:
                file_name = chunk["source"].get("file", "unknown file")
                page = chunk["source"].get("page")
                if page:
                    page = str(page).strip()
                    m = re.search(r'\d+', page)
                    page = m.group(0) if m else "N/A"
                else:
                    page = "N/A"
                
                results.append({
                    "file": file_name,
                    "page": page,
                    "snippets": snippets,
                    "score": match_score,
                    "exact_match": exact_match
                })
    
    # Sort by score (descending), exact match status, and then file/page
    return sorted(results, key=lambda x: (-x["score"], -int(x["exact_match"]), x["file"].lower(), int(x["page"]) if x["page"].isdigit() else float('inf')))

# -----------------------------
# File Upload & Processing
# -----------------------------
uploaded_files = st.file_uploader(
    "Upload one or more files or a ZIP file containing a folder (.txt, .pdf, .csv, .pptx, .json, .zip):",
    type=["txt", "pdf", "csv", "pptx", "json", "zip"],
    accept_multiple_files=True
)

if uploaded_files and "embeddings" not in st.session_state:
    if "chunks" not in st.session_state:
        st.session_state.chunks = []
    
    for uploaded_file in uploaded_files:
        if uploaded_file.size > MAX_FILE_SIZE:
            st.error(f"File {uploaded_file.name} exceeds the size limit of {MAX_FILE_SIZE / (1024 * 1024):.1f} MB.")
            continue

        if uploaded_file.name.lower().endswith(".zip"):
            with st.spinner(f"Extracting folder from {uploaded_file.name}..."):
                try:
                    with zipfile.ZipFile(uploaded_file, 'r') as zip_ref:
                        for zip_info in zip_ref.infolist():
                            if not zip_info.is_dir():
                                ext = zip_info.filename.split('.')[-1].lower()
                                if ext in ["txt", "pdf", "csv", "pptx", "json"]:
                                    file_content = zip_ref.read(zip_info.filename)
                                    file_bytes = BytesIO(file_content)
                                    chunks = process_file_bytes(file_bytes, zip_info.filename)
                                    if chunks:
                                        st.session_state.chunks.extend(chunks)
                                        st.success(f"Processed file from folder: {zip_info.filename}")
                                else:
                                    st.warning(f"Skipping unsupported file {zip_info.filename} in the ZIP.")
                except Exception as e:
                    st.error(f"Failed to extract ZIP file {uploaded_file.name}: {str(e)}")
        else:
            with st.spinner(f"Processing file {uploaded_file.name}..."):
                chunks = process_file(uploaded_file)
                if chunks:
                    st.session_state.chunks.extend(chunks)
                    st.success(f"Processed file: {uploaded_file.name}")
    
    text_for_embedding = [chunk["text"] for chunk in st.session_state.chunks]
    with st.spinner("Generating embeddings..."):
        embeddings = generate_embeddings(text_for_embedding)
        if embeddings.size > 0:
            st.session_state.embeddings = embeddings
            st.success("Embeddings generated successfully!")
        else:
            st.error("Failed to generate embeddings. Please try again.")

# -----------------------------
# Chat Interface
# -----------------------------
if "messages" not in st.session_state:
    st.session_state.messages = []

for message in st.session_state.messages:
    with st.chat_message(message["role"]):
        st.markdown(message["content"])

if not st.session_state.get("chunks"):
    st.warning("Please upload one or more files first.")
    st.stop()

if prompt := st.chat_input("Ask a question about the uploaded content:"):
    st.session_state.messages.append({"role": "user", "content": prompt})
    with st.chat_message("user"):
        st.markdown(prompt)

    # Patterns to identify search queries for any term
    patterns = [
        r'(?:exact matches?|find|show|get|list|count|how many|mentions|occurrences) (?:of )?(?:the\s+)?(?:word\s+|term\s+)?[\'"]?([^\'"\n.,;]+?)[\'"]?(?:\s|$)',
        r'(?:search|look) for [\'"]?([^\'"\n.,;]+?)[\'"]?(?:\s|$)',
        r'(?:the\s+)?(?:word|term)\s+[\'"]?([^\'"\n.,;]+?)[\'"]?',
        r'occurrences? of [\'"]?([^\'"\n.,;]+?)[\'"]?(?:\s|$)'
    ]
    
    # Try to identify search term in different patterns
    search_term = None
    is_search_query = False
    
    # First check if this is a search query
    for pattern in patterns:
        if re.search(pattern, prompt, re.IGNORECASE):
            is_search_query = True
            break
    
    if is_search_query:
        # Extract the actual search term
        for pattern in patterns:
            search_match = re.search(pattern, prompt, re.IGNORECASE)
            if search_match:
                term = search_match.group(1).strip()
                # Clean and normalize search term
                term = re.sub(r'[.,;:]', '', term).strip()
                if len(term) > 2 and term.lower() not in ['all', 'the', 'and', 'for', 'but', 'its', 'any', 'are', 'was', 'were']:
                    search_term = term
                    break
        
        if not search_term:
            st.warning("Please specify a valid search term. For example:\n" + 
                      "- 'find Chamber'\n" +
                      "- 'show mentions of [specific word]'\n" +
                      "- 'how many times [word] appears'")

    if search_term:
        exact_results = extract_exact_mentions(st.session_state.chunks, search_term)
        if exact_results:
            total_mentions = 0
            response_lines = []
            for idx, res in enumerate(exact_results, start=1):
                snippets = res["snippets"]
                file_info = f"{idx}. File: {res['file']}, Page: {res['page']}"
                if res["exact_match"]:
                    file_info += " (Exact Match)"
                response_lines.append(f"{file_info}")
                
                for snippet in snippets:
                    total_mentions += len(list(re.finditer(rf'\b{re.escape(search_term)}\b', str(snippet), re.IGNORECASE)))
                    response_lines.append(f"   • {snippet}")
            
            summary = f"Found {total_mentions} total mention(s) of '{search_term}' across {len(exact_results)} sections:\n\n"
            bot_response = summary + "\n\n".join(response_lines)
        else:
            bot_response = f"No occurrences of '{search_term}' were found in the documents. Please check:\n" + \
                          "1. The word is spelled correctly\n" + \
                          "2. The files have been properly uploaded\n" + \
                          "3. Try searching for part of the word or a related term"
        st.session_state.messages.append({"role": "assistant", "content": bot_response})
        with st.chat_message("assistant"):
            st.markdown(bot_response)
    else:
        file_match = re.search(r'file\s+([A-Za-z0-9\-_]+)', prompt, re.IGNORECASE)
        if file_match:
            file_id = file_match.group(1)
            filtered_indices = [
                i for i, chunk in enumerate(st.session_state.chunks)
                if file_id.lower() in chunk["source"].get("file", "").lower()
            ]
            if filtered_indices:
                filtered_chunks = [st.session_state.chunks[i] for i in filtered_indices]
                filtered_text_chunks = [chunk["text"] for chunk in filtered_chunks]
                filtered_embeddings = st.session_state.embeddings[filtered_indices]
            else:
                filtered_chunks = st.session_state.chunks
                filtered_text_chunks = [chunk["text"] for chunk in st.session_state.chunks]
                filtered_embeddings = st.session_state.embeddings
        else:
            filtered_chunks = st.session_state.chunks
            filtered_text_chunks = [chunk["text"] for chunk in st.session_state.chunks]
            filtered_embeddings = st.session_state.embeddings

        with st.spinner("Retrieving relevant context..."):
            top_indices = find_relevant_context_indices(prompt, filtered_text_chunks, filtered_embeddings)
            
        if not top_indices:
            st.error("Could not find relevant context. Please try rephrasing your question.")
            st.session_state.messages.append({"role": "assistant", "content": "Could not find relevant context. Please try rephrasing your question."})
            with st.chat_message("assistant"):
                st.markdown("Could not find relevant context. Please try rephrasing your question.")
        else:
            relevant_chunks = [filtered_chunks[i] for i in top_indices]
            formatted_contexts = []
            
            for chunk in relevant_chunks:
                file_ref = chunk["source"].get("file", "unknown file")
                page_ref = chunk["source"].get("page")
                if page_ref and str(page_ref).strip().isdigit():
                    ref_str = f"File: {file_ref}, Page: {page_ref}"
                else:
                    ref_str = f"File: {file_ref}"
                formatted_context = f"{ref_str}\n{chunk['text']}"
                formatted_contexts.append(formatted_context)
                
            combined_context = "\n\n".join(formatted_contexts)

            with st.spinner("Generating response..."):
                messages = [
                    {"role": "system", "content": "You are a helpful assistant. Always cite the specific file and page numbers when referring to information from the documents."},
                    {"role": "user", "content": f"Context:\n{combined_context}\n\nQuestion: {prompt}"}
                ]
                bot_response = get_chat_response(messages)
                if bot_response:
                    st.session_state.messages.append({"role": "assistant", "content": bot_response})
                    with st.chat_message("assistant"):
                        st.markdown(bot_response)
else:
    st.warning("Please upload file(s) and wait for embeddings to be generated before asking questions.")

